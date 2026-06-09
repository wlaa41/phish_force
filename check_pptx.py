from pptx import Presentation
prs = Presentation(r'C:\Users\Lenovo\Documents\GitHub\OUR Group\ShadowPhish_Team_Presentation.pptx')
print(f'Total slides: {len(prs.slides)}')
for i, s in enumerate(prs.slides, 1):
    texts = []
    for sh in s.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t:
                texts.append(t[:55])
    print(f'  Slide {i}: {" || ".join(texts[:4])}')
