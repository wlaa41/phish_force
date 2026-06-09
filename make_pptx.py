"""
ShadowPhish Cyber Defense Academy - Team Contribution Presentation
12 slides (2 per member x 6 members)
Updated to reflect Microsoft Planner task board (PhishForce)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W, H = Inches(10), Inches(5.625)   # 16:9

# colours
C_BG     = RGBColor(0x07, 0x09, 0x1A)
C_CARD   = RGBColor(0x11, 0x16, 0x30)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LGREY  = RGBColor(0xB0, 0xB8, 0xCC)
C_CYAN   = RGBColor(0x00, 0xC8, 0xFF)
C_RED    = RGBColor(0xFF, 0x50, 0x50)
C_RED_BG = RGBColor(0x1A, 0x0A, 0x0A)
C_RED_TX = RGBColor(0xFF, 0xB0, 0xB0)
C_META   = RGBColor(0x50, 0x58, 0x78)

def rgb_from_hex(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill_rgb, line_rgb=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, x, y, w, h,
                font_size=14, bold=False, italic=False,
                color=None, align=PP_ALIGN.LEFT, font_face="Calibri"):
    if color is None:
        color = C_WHITE
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    run.font.name      = font_face
    return txb

def add_bullets(slide, items, x, y, w, h,
                font_size=12, color=None, font_face="Calibri"):
    if color is None:
        color = C_LGREY
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pPr = p._p.get_or_add_pPr()
        pPr.set('marL', str(int(Emu(Inches(0.12)))))
        pPr.set('indent', str(int(Emu(Inches(-0.12)))))
        if i > 0:
            pPr.set('spcBef', str(int(Pt(2) * 100)))
        run = p.add_run()
        run.text = ">  " + item
        run.font.size      = Pt(font_size)
        run.font.color.rgb = color
        run.font.name      = font_face

def top_stripe(slide, color):
    add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.05), fill_rgb=color)

def project_label(slide):
    add_textbox(slide, "ShadowPhish  -  Coventry University  -  503IT",
                Inches(4.5), Inches(0.10), Inches(5.2), Inches(0.28),
                font_size=9, color=C_META, align=PP_ALIGN.RIGHT)

def name_badge(slide, name, role, accent):
    add_rect(slide, Inches(0.45), Inches(0.55), Inches(0.065), Inches(0.70), fill_rgb=accent)
    add_textbox(slide, name,
                Inches(0.58), Inches(0.52), Inches(6), Inches(0.42),
                font_size=26, bold=True, color=C_WHITE)
    add_textbox(slide, role,
                Inches(0.58), Inches(0.94), Inches(7), Inches(0.30),
                font_size=13, italic=True, color=accent)

def section_tag(slide, label, x, y, accent):
    add_rect(slide, x, y, Inches(2.35), Inches(0.33),
             fill_rgb=C_CARD, line_rgb=accent, line_w=Pt(1.2))
    add_textbox(slide, label, x, y + Inches(0.02), Inches(2.35), Inches(0.30),
                font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

def bullet_card(slide, bullets, x, y, w, h, accent):
    add_rect(slide, x, y, w, h, fill_rgb=C_CARD, line_rgb=accent, line_w=Pt(1.2))
    add_rect(slide, x, y, Inches(0.055), h, fill_rgb=accent)
    add_bullets(slide, bullets,
                x + Inches(0.18), y + Inches(0.12),
                w - Inches(0.28), h - Inches(0.22))

def slide_number(slide, n, total):
    add_textbox(slide, "{} / {}".format(n, total),
                Inches(8.8), Inches(5.20), Inches(0.9), Inches(0.30),
                font_size=10, color=C_META, align=PP_ALIGN.RIGHT)

def challenge_solution_slide(slide, challenge, solution, accent):
    section_tag(slide, "CHALLENGE", Inches(0.45), Inches(1.45), C_RED)
    add_rect(slide, Inches(0.45), Inches(1.85), Inches(9.1), Inches(1.44),
             fill_rgb=C_RED_BG, line_rgb=C_RED, line_w=Pt(1.2))
    add_rect(slide, Inches(0.45), Inches(1.85), Inches(0.055), Inches(1.44), fill_rgb=C_RED)
    add_textbox(slide, challenge,
                Inches(0.65), Inches(1.96), Inches(8.70), Inches(1.22),
                font_size=13, color=C_RED_TX)

    section_tag(slide, "SOLUTION", Inches(0.45), Inches(3.44), accent)
    add_rect(slide, Inches(0.45), Inches(3.84), Inches(9.1), Inches(1.25),
             fill_rgb=C_CARD, line_rgb=accent, line_w=Pt(1.2))
    add_rect(slide, Inches(0.45), Inches(3.84), Inches(0.055), Inches(1.25), fill_rgb=accent)
    add_textbox(slide, solution,
                Inches(0.65), Inches(3.95), Inches(8.70), Inches(1.05),
                font_size=13, color=C_LGREY)


# =======================================================================
# TEAM DATA  (Planner board: PhishForce)
# Buckets: backlogs | To do | doing | done | Review
# =======================================================================

members = [
    dict(
        name="Prateek Rana",
        role="Project Lead  -  Backend, Architecture & Story",
        hex="8B00FF",
        left_title="KEY CONTRIBUTIONS",
        left_bullets=[
            "Defined Game Design & Architecture - the blueprint for the project",
            "Completed Backend API Setup - core server logic built and merged",
            "Built Frontend UI Framework used across all game pages",
            "Designed Graphics & Icons for the cyberpunk visual theme",
            "Wrote the game Story and led Documentation with Prathish",
            "Planned Trophy System and Daily Challenge System features",
            "Kept Planner board updated throughout the project lifecycle",
        ],
        right_title="PLANNER TASKS",
        right_bullets=[
            "Game Design & Architecture  [done]",
            "Backend API Setup  [done - Completed]",
            "Frontend UI Framework  [done]",
            "Graphics & Icons  [done]",
            "Documentation  [done]",
            "Story Writing  [Review - in sign-off]",
            "Trophy System  [To do]",
            "Daily Challenge System  [done]",
            "QA Testing - full game  [done]",
        ],
        challenge=(
            "The Intel Mission feedback bar shifted the entire quiz layout every time "
            "a player submitted an answer, making the game feel broken and pulling "
            "attention away from the phishing awareness content."
        ),
        solution=(
            "Replaced display:none / display:block with visibility:hidden / visible "
            "and added a fixed min-height:56px to the bar element. The space is always "
            "reserved so no reflow occurs on answer - a one-rule CSS fix that eliminated "
            "all layout jumping during gameplay."
        ),
    ),
    dict(
        name="Prathish Thangaraj",
        role="UI Developer  -  Screen Design & Load System",
        hex="00C8FF",
        left_title="KEY CONTRIBUTIONS",
        left_bullets=[
            "Completed Load Screen - first major screen fully built and merged",
            "Designing Level Screen layout (currently in progress)",
            "Co-designed Database Structure with Allen for player data model",
            "Co-authored Game Design & Architecture and Documentation with Prateek",
            "Proposed Dark Mode UI feature for visual accessibility (in backlog)",
            "Developing Trophy System reward mechanics with Prateek",
        ],
        right_title="PLANNER TASKS",
        right_bullets=[
            "Load Screen  [done - Completed]",
            "Level Screen Design  [doing - In progress]",
            "Database Structure Design  [done]",
            "Game Design & Architecture  [done]",
            "Documentation  [done]",
            "Trophy System  [To do]",
            "Dark Mode UI  [backlogs]",
        ],
        challenge=(
            "Player progress (coins, hearts, XP) was being silently lost whenever the "
            "user navigated between pages. The game appeared to save correctly but "
            "nothing persisted after a page reload."
        ),
        solution=(
            "Traced the bug to sessionStorage being used in the state initialisation "
            "function. sessionStorage is wiped when the tab closes, unlike localStorage. "
            "Replacing that single API call fixed persistence across all pages and "
            "browser sessions with no further changes needed."
        ),
    ),
    dict(
        name="Allen Twins George",
        role="UI/UX Designer  -  Accounts & Sound Integration",
        hex="00E676",
        left_title="KEY CONTRIBUTIONS",
        left_bullets=[
            "Built the Account System - login, registration and guest access flow",
            "Co-designed Database Structure with Prathish for player data model",
            "Working on Sound Integration to add audio feedback to gameplay",
            "Co-developed Daily Challenge System with Prateek",
            "Created UI wireframes and defined the game's visual identity",
        ],
        right_title="PLANNER TASKS",
        right_bullets=[
            "Account System  [done]",
            "Database Structure Design  [done]",
            "Sound Integration  [doing - In progress]",
            "Daily Challenge System  [done]",
        ],
        challenge=(
            "Road obstacles in Highway Chase appeared suddenly in the middle of the "
            "screen, looking like a glitch and giving players zero reaction time - "
            "the game felt completely unfair."
        ),
        solution=(
            "Changed the obstacle spawn Y coordinate from H x 0.38 - 80 (mid-canvas) "
            "to -90 (above the canvas edge). Objects now enter from off-screen and "
            "scroll naturally into view, giving players a clear approach window. "
            "A single number change transformed the fairness and feel of the game."
        ),
    ),
    dict(
        name="Bishal Raut",
        role="Game Logic  -  Economy, Hints & Database",
        hex="FFB300",
        left_title="KEY CONTRIBUTIONS",
        left_bullets=[
            "Implemented Hint System - hints cost coins and reveal partial answers",
            "Balanced Coin & Heart Economy across all game modes",
            "Building Database Player Progress tracker (currently in progress)",
            "Co-built Core Game Logic with Avtar powering the main gameplay loop",
            "Participated in full-game QA Testing pass with the whole team",
        ],
        right_title="PLANNER TASKS",
        right_bullets=[
            "Hint System  [done]",
            "Coin & Heart Economy Balancing  [done]",
            "Database Player Progress  [doing - In progress]",
            "Core Game Logic  [done]",
            "QA Testing - full game  [done]",
        ],
        challenge=(
            "The original question bank had only 5 questions recycled randomly, making "
            "Intel Mission feel repetitive. Players who replayed saw the same questions "
            "immediately, removing all educational value on repeat sessions."
        ),
        solution=(
            "Expanded the bank to 20 unique real-world phishing scenarios covering the "
            "full taxonomy from module 503IT - spot-the-phish, URL analysis, response-"
            "choice, classify, and fill-in-the-blank. Questions draw in randomised order "
            "so every session stays fresh and academically rigorous."
        ),
    ),
    dict(
        name="Dipesh Regmi",
        role="Frontend Developer  -  UI, Leaderboard & Player Data",
        hex="FF6E40",
        left_title="KEY CONTRIBUTIONS",
        left_bullets=[
            "Built Leaderboard UI to display and rank top player scores",
            "Implemented Username Management for personalised player profiles",
            "Collaborating on Database Player Progress tracking with Bishal",
            "Working on Sound Integration feature alongside Allen",
            "Verified cross-browser compatibility on Chrome, Firefox and Edge",
        ],
        right_title="PLANNER TASKS",
        right_bullets=[
            "Leaderboard UI  [done]",
            "Username Management  [done]",
            "Database Player Progress  [doing - In progress]",
            "Sound Integration  [doing - In progress]",
        ],
        challenge=(
            "The login form allowed empty submissions and short passwords with no "
            "visible feedback - players thought the game was broken when the login "
            "button silently did nothing on invalid input."
        ),
        solution=(
            "Added explicit client-side validation: empty-field checks, minimum "
            "codename length of 3 and password length of 4, coloured error messages, "
            "and a CSS shake animation on the card. The shake gives immediate tactile "
            "feedback so players know exactly what needs correcting before retrying."
        ),
    ),
    dict(
        name="Avtar",
        role="Game Content  -  Questions, Core Logic & QA",
        hex="E040FB",
        left_title="KEY CONTRIBUTIONS",
        left_bullets=[
            "Developing Question Content for all game levels (currently in progress)",
            "Co-built Core Game Logic with Bishal powering the main gameplay loop",
            "Participated in full-game QA Testing alongside the whole team",
            "Integrated SVG ShadowPhish logo with CSS float animation on home screen",
            "Tested infinite runner mode - confirmed stable past 10,000 m",
        ],
        right_title="PLANNER TASKS",
        right_bullets=[
            "Question Content  [doing - In progress]",
            "Core Game Logic  [done]",
            "QA Testing - full game  [done]",
        ],
        challenge=(
            "The original Highway Chase ended after exactly 3,000 m with a WIN screen, "
            "making the game feel too short and removing any incentive for players to "
            "keep answering questions or practising phishing awareness."
        ),
        solution=(
            "Removed the WIN_DISTANCE = 3000 constant and its win-check from the "
            "update() loop entirely. The distance meter was repurposed to show progress "
            "within the current 1,500 m checkpoint segment - giving rhythm and reward "
            "cues without any artificial finish line."
        ),
    ),
]


# =======================================================================
# BUILD PRESENTATION
# =======================================================================

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank

for idx, m in enumerate(members):
    accent = rgb_from_hex(m["hex"])
    sn_a   = idx * 2 + 1
    sn_b   = idx * 2 + 2

    # SLIDE A: contributions
    sA = prs.slides.add_slide(blank_layout)
    set_bg(sA, C_BG)
    top_stripe(sA, accent)
    project_label(sA)
    name_badge(sA, m["name"], m["role"], accent)

    section_tag(sA, m["left_title"], Inches(0.45), Inches(1.45), accent)
    bullet_card(sA, m["left_bullets"],
                Inches(0.45), Inches(1.85), Inches(4.45), Inches(3.20), accent)

    section_tag(sA, m["right_title"], Inches(5.25), Inches(1.45), C_CYAN)
    bullet_card(sA, m["right_bullets"],
                Inches(5.25), Inches(1.85), Inches(4.30), Inches(3.20), C_CYAN)

    slide_number(sA, sn_a, 12)

    # SLIDE B: challenge / solution
    sB = prs.slides.add_slide(blank_layout)
    set_bg(sB, C_BG)
    top_stripe(sB, accent)
    project_label(sB)
    name_badge(sB, m["name"], "Challenge  &  Solution", accent)
    challenge_solution_slide(sB, m["challenge"], m["solution"], accent)
    slide_number(sB, sn_b, 12)


OUT = r"C:\Users\Lenovo\Documents\GitHub\OUR Group\ShadowPhish_Team_Presentation.pptx"
prs.save(OUT)
print("Saved: " + OUT)
